#!/usr/bin/env python3
"""
建設業特化AI顧問サービス プレゼンテーション資料
NOVALIS デザインテンプレート準拠
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# =============================================================================
# デザイン定数
# =============================================================================
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# カラーパレット
COLOR_PINK = RGBColor(237, 30, 121)       # #ED1E79 - アクセント
COLOR_SUB_PINK = RGBColor(245, 105, 117)  # #F56975 - グラデーション用
COLOR_BLACK = RGBColor(0, 0, 0)           # #000000
COLOR_WHITE = RGBColor(255, 255, 255)     # #FFFFFF
COLOR_BG_GRAY = RGBColor(245, 248, 248)   # #f5f8f8 - 背景
COLOR_TEXT_GRAY = RGBColor(155, 155, 155) # #9B9B9B
COLOR_DARK_GRAY = RGBColor(74, 74, 74)    # #4A4A4A
COLOR_SECTION_NUM = RGBColor(102, 102, 102)  # #666666

# フォント
FONT_EN = "Oswald"
FONT_JP = "Noto Sans JP"

# レイアウト寸法
HEADER_HEIGHT = Inches(0.9)  # 約12%
PINK_BAR_WIDTH = Inches(0.15)
MARGIN = Inches(0.5)

# ロゴパス
LOGO_PATH = "/Users/tanakashunsuke/.claude/skills/novalis-slide-template/assets/NOVALIS3.png"

# =============================================================================
# ヘルパー関数
# =============================================================================

def set_shape_fill(shape, color):
    """図形の塗りつぶし色を設定"""
    shape.fill.solid()
    shape.fill.fore_color.rgb = color

def set_shape_no_line(shape):
    """図形の線を消す"""
    shape.line.fill.background()

def add_text_frame(shape, text, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.LEFT):
    """テキストフレームにテキストを追加"""
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    return tf

def add_paragraph(text_frame, text, font_name, font_size, font_color, bold=False, alignment=PP_ALIGN.LEFT, space_before=Pt(0), space_after=Pt(6)):
    """テキストフレームに段落を追加"""
    p = text_frame.add_paragraph()
    p.alignment = alignment
    p.space_before = space_before
    p.space_after = space_after
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = font_color
    run.font.bold = bold
    return p

def create_content_slide_base(prs, eng_title, page_num=None):
    """コンテンツスライドのベースを作成（ヘッダー、ピンクバー、ロゴ）"""
    slide_layout = prs.slide_layouts[6]  # 空白レイアウト
    slide = prs.slides.add_slide(slide_layout)

    # 背景色設定
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_GRAY

    # 黒ヘッダー
    header = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        SLIDE_WIDTH, HEADER_HEIGHT
    )
    set_shape_fill(header, COLOR_BLACK)
    set_shape_no_line(header)

    # ピンク縦バー
    pink_bar = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        PINK_BAR_WIDTH, HEADER_HEIGHT
    )
    set_shape_fill(pink_bar, COLOR_PINK)
    set_shape_no_line(pink_bar)

    # 英語タイトル
    title_box = slide.shapes.add_textbox(
        Inches(0.4), Inches(0.25),
        Inches(8), Inches(0.5)
    )
    add_text_frame(title_box, eng_title, FONT_EN, Pt(28), COLOR_WHITE, bold=True)

    # ロゴ
    if os.path.exists(LOGO_PATH):
        logo = slide.shapes.add_picture(
            LOGO_PATH,
            SLIDE_WIDTH - Inches(1.8), Inches(0.2),
            height=Inches(0.5)
        )

    # ページ番号
    if page_num:
        page_box = slide.shapes.add_textbox(
            SLIDE_WIDTH / 2 - Inches(0.25), SLIDE_HEIGHT - Inches(0.4),
            Inches(0.5), Inches(0.3)
        )
        add_text_frame(page_box, str(page_num), FONT_EN, Pt(12), COLOR_TEXT_GRAY, alignment=PP_ALIGN.CENTER)

    return slide

def add_white_content_box(slide, left, top, width, height):
    """白いコンテンツボックスを追加"""
    box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        left, top, width, height
    )
    set_shape_fill(box, COLOR_WHITE)
    set_shape_no_line(box)
    # 角丸を小さく
    box.adjustments[0] = 0.02
    return box

# =============================================================================
# スライド作成関数
# =============================================================================

def create_cover_slide(prs):
    """ページ1：表紙スライド"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # 黒背景
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BLACK

    # メインキャッチコピー（日本語）
    catch_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(1.5),
        Inches(11.5), Inches(1.2)
    )
    tf = catch_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "「AIを使いたいけど、何から始めれば...」"
    run.font.name = FONT_JP
    run.font.size = Pt(32)
    run.font.color.rgb = COLOR_TEXT_GRAY

    # 英語タイトル
    eng_title_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.5),
        Inches(11.5), Inches(1.5)
    )
    tf = eng_title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "AI ADVISORY SERVICE"
    run.font.name = FONT_EN
    run.font.size = Pt(72)
    run.font.color.rgb = COLOR_WHITE
    run.font.bold = True

    # サブタイトル
    sub_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(4.0),
        Inches(11.5), Inches(0.8)
    )
    tf = sub_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "建設業専門のAI顧問が、月10万円で御社に"
    run.font.name = FONT_JP
    run.font.size = Pt(28)
    run.font.color.rgb = COLOR_TEXT_GRAY

    # 実績コピー
    cred_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(5.0),
        Inches(11.5), Inches(1.0)
    )
    tf = cred_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "複数回全国トップセールス獲得・5年連続で個人年間売上2億円維持\n建設業の現場を知り尽くしたAI専門家が、御社のAI活用を0から伴走支援"
    run.font.name = FONT_JP
    run.font.size = Pt(14)
    run.font.color.rgb = COLOR_TEXT_GRAY

    # ロゴ
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(
            LOGO_PATH,
            SLIDE_WIDTH - Inches(2.0), SLIDE_HEIGHT - Inches(0.8),
            height=Inches(0.5)
        )

    return slide

def create_problem_slide(prs):
    """ページ2：こんなお悩みありませんか？"""
    slide = create_content_slide_base(prs, "Problems", 2)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.6)
    )
    add_text_frame(jp_title, "こんなお悩みありませんか？", FONT_JP, Pt(28), COLOR_BLACK, bold=True)

    # 白いコンテンツボックス
    content_box = add_white_content_box(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.8))

    # 課題リスト用テキストボックス
    problems = [
        "□ 日中は現場、定時後に事務処理。気づけば毎月35時間以上の残業",
        "□ 見積もり作成に2時間以上。原価表を見ながら電卓を叩く日々",
        "□ 日報・写真整理・報告書作成。この『ちょっとした作業』の積み重ねが残業に",
        "□ 提案資料がいつも似たようなものになり、差別化できず契約率が上がらない",
        "□ チラシ制作を外注すると約20万円。自分で作りたいがデザインスキルも時間もない",
        "□ AIを使いたいが、何から始めればいいかわからない。高額な投資のイメージもある",
        "□ AI人材を採用したいが年収も高額。でも若手に教えられるほど自分も詳しくない",
    ]

    problem_box = slide.shapes.add_textbox(
        Inches(0.8), Inches(2.0),
        Inches(11.5), Inches(3.8)
    )
    tf = problem_box.text_frame
    tf.word_wrap = True

    for i, problem in enumerate(problems):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(12)
        run = p.add_run()
        run.text = problem
        run.font.name = FONT_JP
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_BLACK

    # 締めの一言
    closing_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.0),
        Inches(12.3), Inches(1.0)
    )
    tf = closing_box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "「AIを導入したいけど、どこから手をつければ...」\nその悩み、建設業で毎日0時残業から定時帰りを実現した私が、0から一緒に解決します。"
    run.font.name = FONT_JP
    run.font.size = Pt(14)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    return slide

def create_why_fail_slide(prs):
    """ページ3：なぜ多くの会社がAI導入に失敗するのか"""
    slide = create_content_slide_base(prs, "Why AI Projects Fail", 3)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.6)
    )
    add_text_frame(jp_title, "なぜ多くの会社がAI導入に失敗するのか", FONT_JP, Pt(26), COLOR_BLACK, bold=True)

    # 左側コンテンツボックス（多くの人が想像する効率化）
    left_box = add_white_content_box(slide, Inches(0.5), Inches(1.8), Inches(5.8), Inches(2.2))

    left_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.9),
        Inches(5.4), Inches(2.0)
    )
    tf = left_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "多くの人が想像する『効率化』"
    run.font.name = FONT_JP
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = "「見積もり作成が、ボタン一つで終わる」\n「提案書が、自動で完璧に仕上がる」\n\nたしかに、AIがあれば実現可能です。\nしかし、最大の効率化とは、\nもっと地味な改善の積み重ねです。"
    run.font.name = FONT_JP
    run.font.size = Pt(13)
    run.font.color.rgb = COLOR_BLACK

    # 右側コンテンツボックス（本当の効率化）
    right_box = add_white_content_box(slide, Inches(6.5), Inches(1.8), Inches(6.3), Inches(2.2))

    right_text = slide.shapes.add_textbox(
        Inches(6.7), Inches(1.9),
        Inches(5.9), Inches(2.0)
    )
    tf = right_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "本当の効率化とは"
    run.font.name = FONT_JP
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    p = tf.add_paragraph()
    p.space_before = Pt(8)
    run = p.add_run()
    run.text = "「原価を調べる5分」を2分に。\n「文章を考える3分」を1分に。\n「ファイル名をつける2分」を30秒に。\n\n5分の短縮を10個実現するだけで、50分。\nこれを10日やったら、500分。"
    run.font.name = FONT_JP
    run.font.size = Pt(13)
    run.font.color.rgb = COLOR_BLACK

    # 下部コンテンツボックス（だから「AI顧問」）
    bottom_box = add_white_content_box(slide, Inches(0.5), Inches(4.2), Inches(12.3), Inches(2.4))

    bottom_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(4.3),
        Inches(11.9), Inches(2.2)
    )
    tf = bottom_text.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "でも、アプリや外注では解決しない"
    run.font.name = FONT_JP
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    p = tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = "「常に、あなたの会社のどこを効率化できるか」を見極め続ける人が必要だから。\n業務はどんなものがあって、どう分解すればいいのか。どこにAIが使えて、どこに使えないのか。\nそれを判断して、解決策まで導く。これは、専門家がいないとできません。"
    run.font.name = FONT_JP
    run.font.size = Pt(13)
    run.font.color.rgb = COLOR_BLACK

    p = tf.add_paragraph()
    p.space_before = Pt(12)
    run = p.add_run()
    run.text = "だから、「AI顧問」という形を作りました。常に寄り添ってくれる人。伴走してくれる人。それが、このサービスの本質です。"
    run.font.name = FONT_JP
    run.font.size = Pt(14)
    run.font.color.rgb = COLOR_BLACK
    run.font.bold = True

    return slide

def create_story_slide(prs):
    """ページ4：なぜ私が建設業に特化するのか（原体験ストーリー）"""
    slide = create_content_slide_base(prs, "My Story", 4)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.6)
    )
    add_text_frame(jp_title, "なぜ私が建設業に特化するのか", FONT_JP, Pt(26), COLOR_BLACK, bold=True)

    # 左側：ストーリーボックス
    story_box = add_white_content_box(slide, Inches(0.5), Inches(1.8), Inches(7.5), Inches(5.0))

    story_text = slide.shapes.add_textbox(
        Inches(0.7), Inches(1.95),
        Inches(7.1), Inches(4.7)
    )
    tf = story_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "私がこのサービスを作った理由"
    run.font.name = FONT_JP
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    story_content = """
リフォーム営業として入社した頃、毎日0時を超える残業が当たり前でした。

見積もり作成に2時間以上。原価表すらなく、FAXを指で照らし合わせて原価を調べる。日中は現場、帰社してから事務作業。どれだけ頑張っても、12時を切ることがない。

ある時、気づきました。
「こんな環境で、新しい社員が定着するわけがない。」

だから私は、効率化を始めました。

提案資料のテンプレートを作り直し、原価表を自分で更新し、見積もりシステムを自作。社内に数え切れないほどあった報告書式。目につくものを片っ端から直していきました。

身の回りのすべてを、端から端まで効率化しました。"""

    p = tf.add_paragraph()
    p.space_before = Pt(6)
    run = p.add_run()
    run.text = story_content.strip()
    run.font.name = FONT_JP
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_BLACK

    # 右側：実績ボックス
    result_box = add_white_content_box(slide, Inches(8.2), Inches(1.8), Inches(4.6), Inches(5.0))

    result_text = slide.shapes.add_textbox(
        Inches(8.4), Inches(1.95),
        Inches(4.2), Inches(4.7)
    )
    tf = result_text.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "その結果"
    run.font.name = FONT_JP
    run.font.size = Pt(16)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    results = [
        ("見積もり作成", "2時間 → 10分"),
        ("毎日0時残業", "→ 定時帰り"),
        ("全国トップセールス", "複数回獲得"),
        ("5年連続", "年間売上2億円"),
        ("創業40年で", "過去最高売上達成"),
    ]

    for label, value in results:
        p = tf.add_paragraph()
        p.space_before = Pt(16)
        run = p.add_run()
        run.text = label
        run.font.name = FONT_JP
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_TEXT_GRAY

        p = tf.add_paragraph()
        p.space_before = Pt(2)
        run = p.add_run()
        run.text = value
        run.font.name = FONT_JP
        run.font.size = Pt(18)
        run.font.color.rgb = COLOR_BLACK
        run.font.bold = True

    return slide

def create_concept_slide(prs):
    """ページ5：サービスのコンセプト"""
    slide = create_content_slide_base(prs, "Service Concept", 5)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.6)
    )
    add_text_frame(jp_title, "AI人材を「採用」するのではなく、「顧問」として迎える", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    # 比較表ボックス
    table_box = add_white_content_box(slide, Inches(0.5), Inches(1.8), Inches(12.3), Inches(4.2))

    # 表ヘッダー
    # 左列ヘッダー
    header_left = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0.7), Inches(2.0),
        Inches(5.8), Inches(0.6)
    )
    set_shape_fill(header_left, COLOR_BLACK)
    set_shape_no_line(header_left)
    header_left_text = slide.shapes.add_textbox(Inches(0.7), Inches(2.05), Inches(5.8), Inches(0.5))
    add_text_frame(header_left_text, "AI人材を採用", FONT_JP, Pt(18), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 右列ヘッダー
    header_right = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(6.7), Inches(2.0),
        Inches(5.8), Inches(0.6)
    )
    set_shape_fill(header_right, COLOR_PINK)
    set_shape_no_line(header_right)
    header_right_text = slide.shapes.add_textbox(Inches(6.7), Inches(2.05), Inches(5.8), Inches(0.5))
    add_text_frame(header_right_text, "AI顧問", FONT_JP, Pt(18), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # 比較内容
    comparisons = [
        ("月額 25万円〜", "月額 10万円"),
        ("採用活動が必要（求人・面接・選考）", "今すぐ始められる"),
        ("教育しなければならない", "教育してくれる"),
        ("スキルの見極めが困難", "どんなAIにも対応可能"),
        ("辞めるリスクがある", "辞めない"),
    ]

    y_start = 2.8
    for i, (left_text, right_text) in enumerate(comparisons):
        y = y_start + i * 0.65

        # 左列
        left_box = slide.shapes.add_textbox(Inches(0.7), Inches(y), Inches(5.8), Inches(0.55))
        tf = left_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = left_text
        run.font.name = FONT_JP
        run.font.size = Pt(16) if i == 0 else Pt(14)
        run.font.color.rgb = COLOR_BLACK
        run.font.bold = (i == 0)

        # 右列
        right_box = slide.shapes.add_textbox(Inches(6.7), Inches(y), Inches(5.8), Inches(0.55))
        tf = right_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = right_text
        run.font.name = FONT_JP
        run.font.size = Pt(16) if i == 0 else Pt(14)
        run.font.color.rgb = COLOR_PINK if i == 0 else COLOR_BLACK
        run.font.bold = (i == 0)

    # キーメッセージ
    key_msg = slide.shapes.add_textbox(
        Inches(0.5), Inches(6.2),
        Inches(12.3), Inches(0.8)
    )
    tf = key_msg.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "「この人に聞けば、AI周りはなんとかなる」そんな安心感を、月10万円で。"
    run.font.name = FONT_JP
    run.font.size = Pt(18)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    return slide

def create_plan_slide(prs):
    """ページ6：3つのプラン"""
    slide = create_content_slide_base(prs, "Pricing Plans", 6)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.6)
    )
    add_text_frame(jp_title, "3つのプランからお選びいただけます", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    plans = [
        {
            "name": "伴走プラン",
            "price": "月額 15万円",
            "catch": "「AI専門家が、御社のそばに」",
            "features": ["チャット無制限", "キックオフMTG", "振り返りMTG", "月次レポート"],
            "for": "まずはAI活用を始めたい方向け",
            "highlight": False
        },
        {
            "name": "自走プラン",
            "price": "月額 40万円",
            "catch": "「社内にAI人材を育てる」",
            "features": ["伴走プラン全内容", "社員研修（4名まで）", "内製化支援"],
            "for": "社員にAIスキルを身につけさせたい方向け",
            "highlight": True
        },
        {
            "name": "エージェント開発プラン",
            "price": "月額 60万円",
            "catch": "「御社専用のAIツールを開発」",
            "features": ["要件整理MTG", "月1開発MTG", "オーダーメイド開発"],
            "for": "「これを作ってほしい」がある方向け",
            "highlight": False
        }
    ]

    box_width = Inches(3.9)
    box_height = Inches(5.0)
    start_x = Inches(0.5)
    gap = Inches(0.2)

    for i, plan in enumerate(plans):
        x = start_x + i * (box_width + gap)

        # プランボックス
        plan_box = add_white_content_box(slide, x, Inches(1.8), box_width, box_height)

        # プラン名ヘッダー
        header_color = COLOR_PINK if plan["highlight"] else COLOR_BLACK
        header = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x + Inches(0.1), Inches(1.9),
            box_width - Inches(0.2), Inches(0.5)
        )
        set_shape_fill(header, header_color)
        set_shape_no_line(header)

        header_text = slide.shapes.add_textbox(x + Inches(0.1), Inches(1.92), box_width - Inches(0.2), Inches(0.45))
        add_text_frame(header_text, plan["name"], FONT_JP, Pt(18), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # 価格
        price_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(2.5), box_width - Inches(0.2), Inches(0.5))
        tf = price_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = plan["price"]
        run.font.name = FONT_JP
        run.font.size = Pt(22)
        run.font.color.rgb = COLOR_PINK if plan["highlight"] else COLOR_BLACK
        run.font.bold = True

        # キャッチ
        catch_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(3.0), box_width - Inches(0.2), Inches(0.4))
        tf = catch_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = plan["catch"]
        run.font.name = FONT_JP
        run.font.size = Pt(11)
        run.font.color.rgb = COLOR_TEXT_GRAY

        # 特徴リスト
        features_box = slide.shapes.add_textbox(x + Inches(0.3), Inches(3.5), box_width - Inches(0.4), Inches(2.2))
        tf = features_box.text_frame
        tf.word_wrap = True

        for j, feature in enumerate(plan["features"]):
            if j == 0:
                p = tf.paragraphs[0]
            else:
                p = tf.add_paragraph()
            p.space_after = Pt(8)
            run = p.add_run()
            run.text = f"・{feature}"
            run.font.name = FONT_JP
            run.font.size = Pt(13)
            run.font.color.rgb = COLOR_BLACK

        # 対象者
        for_box = slide.shapes.add_textbox(x + Inches(0.1), Inches(6.3), box_width - Inches(0.2), Inches(0.5))
        tf = for_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = plan["for"]
        run.font.name = FONT_JP
        run.font.size = Pt(10)
        run.font.color.rgb = COLOR_TEXT_GRAY

    return slide

def create_comparison_slide(prs):
    """ページ7：プラン比較表"""
    slide = create_content_slide_base(prs, "Plan Comparison", 7)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.5)
    )
    add_text_frame(jp_title, "プラン比較表", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    # 白いコンテンツボックス
    table_bg = add_white_content_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))

    # テーブル構造
    headers = ["項目", "伴走プラン", "自走プラン", "エージェント開発"]
    rows = [
        ["月額", "15万円", "40万円", "60万円"],
        ["チャット相談", "◎ 無制限", "◎ 無制限", "○ 開発関連"],
        ["キックオフMTG", "◎", "◎", "◎"],
        ["振り返りMTG（3ヶ月後）", "◎", "◎", "◎"],
        ["月次レポート", "◎", "◎", "○"],
        ["月1開発MTG", "−", "−", "◎"],
        ["社員研修（4名まで）", "−", "◎", "−"],
        ["内製化支援", "−", "◎", "−"],
        ["オーダーメイド開発", "−", "−", "◎"],
    ]

    col_widths = [Inches(3.5), Inches(2.8), Inches(2.8), Inches(3.0)]
    row_height = Inches(0.48)
    start_x = Inches(0.6)
    start_y = Inches(1.85)

    # ヘッダー行
    x = start_x
    for j, (header, width) in enumerate(zip(headers, col_widths)):
        cell_bg = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            x, start_y,
            width, row_height
        )
        set_shape_fill(cell_bg, COLOR_BLACK)
        set_shape_no_line(cell_bg)

        cell_text = slide.shapes.add_textbox(x, start_y + Inches(0.08), width, row_height - Inches(0.1))
        tf = cell_text.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = header
        run.font.name = FONT_JP
        run.font.size = Pt(13)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = True

        x += width

    # データ行
    for i, row in enumerate(rows):
        y = start_y + (i + 1) * row_height
        x = start_x

        for j, (cell, width) in enumerate(zip(row, col_widths)):
            # 背景（交互色）
            bg_color = COLOR_WHITE if i % 2 == 0 else COLOR_BG_GRAY
            cell_bg = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                x, y,
                width, row_height
            )
            set_shape_fill(cell_bg, bg_color)
            cell_bg.line.color.rgb = RGBColor(224, 224, 224)
            cell_bg.line.width = Pt(0.5)

            # テキスト
            cell_text = slide.shapes.add_textbox(x + Inches(0.05), y + Inches(0.08), width - Inches(0.1), row_height - Inches(0.1))
            tf = cell_text.text_frame
            p = tf.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT
            run = p.add_run()
            run.text = cell
            run.font.name = FONT_JP
            run.font.size = Pt(12)

            # 価格行は強調
            if i == 0 and j > 0:
                run.font.color.rgb = COLOR_PINK
                run.font.bold = True
            elif cell == "◎":
                run.font.color.rgb = COLOR_PINK
            elif cell == "−":
                run.font.color.rgb = COLOR_TEXT_GRAY
            else:
                run.font.color.rgb = COLOR_BLACK

            x += width

    return slide

def create_contract_slide(prs):
    """ページ8：契約条件・ご利用の流れ"""
    slide = create_content_slide_base(prs, "Contract & Flow", 8)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.5)
    )
    add_text_frame(jp_title, "契約条件・ご利用の流れ", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    # 左側：契約条件
    left_box = add_white_content_box(slide, Inches(0.5), Inches(1.7), Inches(4.5), Inches(2.5))

    contract_title = slide.shapes.add_textbox(Inches(0.7), Inches(1.85), Inches(4.1), Inches(0.5))
    add_text_frame(contract_title, "契約条件", FONT_JP, Pt(18), COLOR_PINK, bold=True)

    conditions = [
        ("最低契約期間", "3ヶ月"),
        ("試用期間", "最初の1ヶ月"),
        ("支払い", "月額・前払い"),
    ]

    conditions_box = slide.shapes.add_textbox(Inches(0.7), Inches(2.4), Inches(4.1), Inches(1.6))
    tf = conditions_box.text_frame
    tf.word_wrap = True

    for i, (label, value) in enumerate(conditions):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(10)
        run = p.add_run()
        run.text = f"{label}："
        run.font.name = FONT_JP
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_TEXT_GRAY

        run = p.add_run()
        run.text = value
        run.font.name = FONT_JP
        run.font.size = Pt(14)
        run.font.color.rgb = COLOR_BLACK
        run.font.bold = True

    # 右側：ご利用の流れ
    right_box = add_white_content_box(slide, Inches(5.2), Inches(1.7), Inches(7.6), Inches(5.2))

    flow_title = slide.shapes.add_textbox(Inches(5.4), Inches(1.85), Inches(7.2), Inches(0.5))
    add_text_frame(flow_title, "ご利用の流れ", FONT_JP, Pt(18), COLOR_PINK, bold=True)

    steps = [
        ("①", "無料AI活用診断（30分・Zoom）", "現状の業務をヒアリング、改善ポイントを洗い出し"),
        ("②", "ロードマップ提示", "「何を」「どの順番で」改善するかを明確化"),
        ("③", "ご契約", ""),
        ("④", "サービス開始", "キックオフMTGから開始、チャット相談スタート"),
    ]

    y = Inches(2.4)
    for num, title, desc in steps:
        # 番号
        num_box = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(5.5), y,
            Inches(0.4), Inches(0.4)
        )
        set_shape_fill(num_box, COLOR_PINK)
        set_shape_no_line(num_box)

        num_text = slide.shapes.add_textbox(Inches(5.5), y + Inches(0.05), Inches(0.4), Inches(0.35))
        add_text_frame(num_text, num, FONT_JP, Pt(12), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

        # タイトル
        step_title = slide.shapes.add_textbox(Inches(6.0), y, Inches(6.5), Inches(0.4))
        add_text_frame(step_title, title, FONT_JP, Pt(14), COLOR_BLACK, bold=True)

        # 説明
        if desc:
            step_desc = slide.shapes.add_textbox(Inches(6.0), y + Inches(0.35), Inches(6.5), Inches(0.5))
            add_text_frame(step_desc, desc, FONT_JP, Pt(11), COLOR_TEXT_GRAY)

        # 矢印線
        if num != "④":
            line_y = y + Inches(0.55) if desc else y + Inches(0.45)
            arrow = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE,
                Inches(5.68), line_y,
                Inches(0.04), Inches(0.5) if desc else Inches(0.3)
            )
            set_shape_fill(arrow, COLOR_TEXT_GRAY)
            set_shape_no_line(arrow)

        y += Inches(1.1) if desc else Inches(0.8)

    return slide

def create_qa_slide(prs):
    """ページ9：よくある質問（Q&A）"""
    slide = create_content_slide_base(prs, "FAQ", 9)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.5)
    )
    add_text_frame(jp_title, "よくある質問", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    # 白いコンテンツボックス
    qa_box = add_white_content_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))

    qas = [
        ("Q. なぜ「3ヶ月」なのですか？", "慣れる。習慣化する。日常に溶け込ませる。そこまで伴走して、初めて「効果が出た」と実感できます。"),
        ("Q. 試用期間とは何ですか？", "最初の1ヶ月で相性を確認。万が一「合わない」と感じた場合は、1ヶ月で終了可能です。"),
        ("Q. 途中でプラン変更はできますか？", "はい。アップグレード・ダウングレードどちらも対応しています。"),
        ("Q. チャット相談はどのくらいで返信がありますか？", "24時間以内に返信いたします。"),
        ("Q. どんな相談ができますか？", "AIに関することなら、どんな相談でも可能です。「こんなこと聞いていいのかな？」もお気軽に。"),
    ]

    y = Inches(1.9)
    for q, a in qas:
        # 質問
        q_box = slide.shapes.add_textbox(Inches(0.7), y, Inches(11.9), Inches(0.4))
        add_text_frame(q_box, q, FONT_JP, Pt(13), COLOR_PINK, bold=True)

        # 回答
        a_box = slide.shapes.add_textbox(Inches(0.7), y + Inches(0.35), Inches(11.9), Inches(0.5))
        tf = a_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = f"→ {a}"
        run.font.name = FONT_JP
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_BLACK

        y += Inches(0.95)

    return slide

def create_cta_slide(prs):
    """ページ10：次のステップ（CTA）"""
    slide = create_content_slide_base(prs, "Next Step", 10)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.5)
    )
    add_text_frame(jp_title, "次のステップ", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    # メインコンテンツボックス
    main_box = add_white_content_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))

    # 左側：メッセージ
    msg_box = slide.shapes.add_textbox(Inches(0.7), Inches(1.9), Inches(6.5), Inches(3.0))
    tf = msg_box.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = "まずは無料診断から"
    run.font.name = FONT_JP
    run.font.size = Pt(22)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    message = """
「何を導入すべきかわからない」
「どこを改善すべきかわからない」

そんな状態で大丈夫です。

大事なのは、「どこを直すべきか」を知ること。
それがわかれば、あとは一つずつ進めるだけです。

改善できること、できないこと。
優先すべきこと、後回しでいいこと。
プロの目で、御社のAI活用ポイントを診断します。"""

    p = tf.add_paragraph()
    p.space_before = Pt(12)
    run = p.add_run()
    run.text = message.strip()
    run.font.name = FONT_JP
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_BLACK

    # 右側：CTA詳細
    cta_box = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(7.5), Inches(2.0),
        Inches(5.0), Inches(4.5)
    )
    set_shape_fill(cta_box, COLOR_BG_GRAY)
    set_shape_no_line(cta_box)
    cta_box.adjustments[0] = 0.03

    cta_title = slide.shapes.add_textbox(Inches(7.7), Inches(2.2), Inches(4.6), Inches(0.6))
    tf = cta_title.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "無料AI活用診断"
    run.font.name = FONT_JP
    run.font.size = Pt(20)
    run.font.color.rgb = COLOR_BLACK
    run.font.bold = True

    cta_sub = slide.shapes.add_textbox(Inches(7.7), Inches(2.7), Inches(4.6), Inches(0.4))
    tf = cta_sub.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "30分・Zoom"
    run.font.name = FONT_JP
    run.font.size = Pt(14)
    run.font.color.rgb = COLOR_PINK
    run.font.bold = True

    cta_details = slide.shapes.add_textbox(Inches(7.7), Inches(3.2), Inches(4.6), Inches(2.5))
    tf = cta_details.text_frame
    tf.word_wrap = True

    details = [
        "・御社の業務をヒアリング",
        "・改善ポイント・優先順位をその場でお伝え",
        "",
        "押し売りは一切ありません。",
        "診断だけでもOK。"
    ]

    for i, detail in enumerate(details):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(6)
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = detail
        run.font.name = FONT_JP
        run.font.size = Pt(12)
        run.font.color.rgb = COLOR_BLACK if i < 3 else COLOR_PINK
        run.font.bold = (i >= 3)

    return slide

def create_contact_slide(prs):
    """ページ11：お問い合わせ"""
    slide = create_content_slide_base(prs, "Contact", 11)

    # 日本語タイトル
    jp_title = slide.shapes.add_textbox(
        Inches(0.5), Inches(1.1),
        Inches(12), Inches(0.5)
    )
    add_text_frame(jp_title, "お問い合わせ", FONT_JP, Pt(24), COLOR_BLACK, bold=True)

    # メインコンテンツボックス
    main_box = add_white_content_box(slide, Inches(0.5), Inches(1.7), Inches(12.3), Inches(5.2))

    # 連絡先情報
    contact_info = [
        ("WEB", "https://novalisgroup.jp/"),
        ("Email", "shunsuke.tanaka@novalisgroup.biz"),
        ("住所", "〒152-0004 東京都目黒区鷹番2丁目20番20号\nイニッゾ学芸大学5-17"),
    ]

    y = Inches(2.3)
    for label, value in contact_info:
        # ラベル
        label_box = slide.shapes.add_textbox(Inches(2.0), y, Inches(1.5), Inches(0.5))
        tf = label_box.text_frame
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.RIGHT
        run = p.add_run()
        run.text = label
        run.font.name = FONT_JP
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_PINK
        run.font.bold = True

        # 値
        value_box = slide.shapes.add_textbox(Inches(3.7), y, Inches(8.0), Inches(0.8))
        tf = value_box.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = value
        run.font.name = FONT_JP
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_BLACK

        y += Inches(1.0) if "\n" not in value else Inches(1.3)

    # 締めのメッセージ
    closing_box = slide.shapes.add_textbox(Inches(0.7), Inches(5.5), Inches(11.9), Inches(1.0))
    tf = closing_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = "お気軽にご相談ください"
    run.font.name = FONT_JP
    run.font.size = Pt(20)
    run.font.color.rgb = COLOR_TEXT_GRAY

    return slide

# =============================================================================
# メイン処理
# =============================================================================

def main():
    """プレゼンテーション作成"""
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    print("スライド作成開始...")

    # 各スライドを作成
    print("  1/11: 表紙")
    create_cover_slide(prs)

    print("  2/11: お悩み")
    create_problem_slide(prs)

    print("  3/11: AI導入失敗の理由")
    create_why_fail_slide(prs)

    print("  4/11: 原体験ストーリー")
    create_story_slide(prs)

    print("  5/11: サービスコンセプト")
    create_concept_slide(prs)

    print("  6/11: 3つのプラン")
    create_plan_slide(prs)

    print("  7/11: プラン比較表")
    create_comparison_slide(prs)

    print("  8/11: 契約条件・ご利用の流れ")
    create_contract_slide(prs)

    print("  9/11: Q&A")
    create_qa_slide(prs)

    print("  10/11: CTA")
    create_cta_slide(prs)

    print("  11/11: お問い合わせ")
    create_contact_slide(prs)

    # 保存
    output_path = "/Users/tanakashunsuke/Desktop/AI-Advisory-Service-Design/AI顧問サービス資料_NOVALIS.pptx"
    prs.save(output_path)
    print(f"\n完成！保存先: {output_path}")

if __name__ == "__main__":
    main()
